from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
import sys
import shlex
import re

# Example input format from file:
"""
template.pptx output.pptx
img headshot.jpg
left=2 top=3 w=3
txt Hello from Python
left=1 top=1 w=6 h=1.5 fontsz=32 align=center
txt Person's name
left=1 top=1 w=6 h=1.5 fontsz=12 align=center
"""

def parse_input(input_file):
    with open(input_file, 'r') as f:
        lines = [line.rstrip('\n') for line in f if line.strip()]
    
    # First line: template.pptx output.pptx
    template, output = lines[0].split()
    blocks = []
    i = 1
    while i < len(lines):
        if lines[i].startswith('img '):
            img_path = lines[i].split()[1]
            img_params = {k: float(v) for k, v in (item.split('=') for item in lines[i+1].split())}
            blocks.append({'type': 'img', 'img_path': img_path, 'params': img_params})
            i += 2
        elif lines[i].startswith('txt '):
            txt = lines[i][4:]
            # Parse combined line with both position and font params
            all_items = shlex.split(lines[i+1])
            txt_params = {}
            font_params = {}
            
            for item in all_items:
                k, v = item.split('=', 1)
                if k in ['left', 'top', 'w', 'h']:
                    txt_params[k] = float(v)
                elif k == 'fontsz':
                    font_params['fontsz'] = int(v)
                elif k == 'align':
                    font_params['align'] = v
                elif k == 'font_cn':
                    # Remove quotes if present
                    if v.startswith('"') and v.endswith('"'):
                        v = v[1:-1]
                    font_params['font_cn'] = v
            
            blocks.append({'type': 'txt', 'txt': txt, 'params': txt_params, 'font_params': font_params})
            i += 2
        else:
            i += 1
    return {
        'template_pptx': template,
        'output_pptx': output,
        'blocks': blocks
    }

def modify_slide(
    template_pptx: str,
    output_pptx: str,
    blocks: list
):
    prs = Presentation(template_pptx)
    slide = prs.slides[0]

    for block in blocks:
        if block['type'] == 'img':
            p = block['params']
            slide.shapes.add_picture(
                block['img_path'],
                Inches(p['left']),
                Inches(p['top']),
                width=Inches(p['w'])
            )
        elif block['type'] == 'txt':
            p = block['params']
            f = block['font_params']
            textbox = slide.shapes.add_textbox(
                Inches(p['left']),
                Inches(p['top']),
                Inches(p['w']),
                Inches(p['h']))
            text_frame = textbox.text_frame
            text_frame.clear()
            para = text_frame.paragraphs[0]

            # Add text
            en, zh = split_en_zh(block['txt'])

            if en:
                r_en = para.add_run()
                r_en.text = en
                r_en.font.name = 'Times New Roman'  # Default English font
                r_en.font.size = Pt(f['fontsz'])

            if zh:
                r_zh = para.add_run()
                r_zh.text = zh
                r_zh.font.size = Pt(f['fontsz'])
                rPr = r_zh.font._element  # this is the <a:rPr> element

                font_cn = f['font_cn'] or "SimSun"

                # Add <a:latin typeface="SimSun"/> if not already present
                latin = rPr.find(qn('a:latin'))
                if latin is None:
                    latin = OxmlElement('a:latin')
                    latin.set('typeface', font_cn)
                    rPr.insert(0, latin)

                # Add <a:ea typeface="SimSun"/>
                ea = rPr.find(qn('a:ea'))
                if ea is None:
                    ea = OxmlElement('a:ea')
                    ea.set('typeface', font_cn)
                    # Insert after <a:latin> if it exists
                    rPr.insert(1, ea)

            if f['align'] == 'center':
                para.alignment = PP_ALIGN.CENTER
            elif f['align'] == 'left':
                para.alignment = PP_ALIGN.LEFT
            elif f['align'] == 'right':
                para.alignment = PP_ALIGN.RIGHT

    prs.save(output_pptx)

CJK_RE = re.compile(
    r'[\u3400-\u4DBF\u4E00-\u9FFF\U00020000-\U0002A6DF]'
)

def split_en_zh(text: str):
    """
    Returns (english_part, chinese_part)
    - english_part may be the full string
    - chinese_part may be empty
    """
    m = CJK_RE.search(text)
    if m is None:
        # English-only (or Latin-only)
        return text, ""

    idx = m.start()
    return text[:idx], text[idx:]

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 script.py <input_file>", file=sys.stderr)
        sys.exit(1)
    
    args = parse_input(sys.argv[1])
    modify_slide(
        template_pptx=args['template_pptx'],
        output_pptx=args['output_pptx'],
        blocks=args['blocks']
    )