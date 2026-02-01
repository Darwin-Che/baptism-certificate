#!/usr/bin/env python3
"""
Combine multiple PPTX files into a single PPTX file.
Usage: python combine_pptx.py output.pptx input1.pptx input2.pptx ...
"""
import sys
from pptx import Presentation
from pptx.util import Inches
from copy import deepcopy
import io

def combine_pptx_files(output_path, input_paths):
    """Combine multiple PPTX files into one by concatenating slides."""
    if not input_paths:
        print("Error: No input files provided", file=sys.stderr)
        sys.exit(1)
    
    # Create a new presentation based on the first file
    combined = Presentation(input_paths[0])
    
    # Add slides from all other presentations
    for input_path in input_paths[1:]:
        try:
            prs = Presentation(input_path)
            for slide in prs.slides:
                # Use blank layout to preserve original formatting
                slide_layout = combined.slide_layouts[6] if len(combined.slide_layouts) > 6 else combined.slide_layouts[0]
                new_slide = combined.slides.add_slide(slide_layout)
                
                # Remove all shapes from the new slide
                for shape in list(new_slide.shapes):
                    sp = shape.element
                    sp.getparent().remove(sp)
                
                # Copy all shapes from source slide, handling images specially
                for shape in slide.shapes:
                    el = shape.element
                    new_el = deepcopy(el)
                    
                    # If this shape contains an image, we need to copy the image relationship
                    if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                        try:
                            # Get the image from the source slide
                            image = shape.image
                            image_bytes = image.blob
                            
                            # Add the image to the new slide's relationships
                            pic = new_slide.shapes.add_picture(
                                io.BytesIO(image_bytes),
                                shape.left,
                                shape.top,
                                shape.width,
                                shape.height
                            )
                            # The image is already added, skip adding the element
                            continue
                        except Exception as e:
                            print(f"Warning: Failed to copy image: {e}", file=sys.stderr)
                    
                    # For non-image shapes, add the copied element
                    new_slide.shapes._spTree.append(new_el)
                    
        except Exception as e:
            print(f"Warning: Failed to add slides from {input_path}: {e}", file=sys.stderr)
            continue
    
    # Save the combined presentation
    combined.save(output_path)
    print(f"Successfully combined {len(input_paths)} presentations with {len(combined.slides)} total slides")

if __name__ == "__main__":
    if len(sys.argv) < 3:
        print("Usage: python combine_pptx.py output.pptx input1.pptx input2.pptx ...", file=sys.stderr)
        sys.exit(1)
    
    output_path = sys.argv[1]
    input_paths = sys.argv[2:]
    
    combine_pptx_files(output_path, input_paths)
